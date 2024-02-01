import docx
import pptx

from ragna.core import DocxDocumentHandler, LocalDocument, PptxDocumentHandler


def get_docx_document(tmp_path, docx_text):
    document = docx.Document()
    document.add_heading(docx_text)
    document.add_paragraph(docx_text)
    path = tmp_path / "test_document.docx"
    document.save(path)
    return LocalDocument.from_path(path)


def test_docx(tmp_path):
    docx_text = "ragna is neat!"
    tmp_docx_document = get_docx_document(tmp_path, docx_text)
    assert isinstance(tmp_docx_document.handler, DocxDocumentHandler)
    pages = list(tmp_docx_document.extract_pages())
    assert len(pages) == 2
    for page in pages:
        assert page.text == docx_text


def get_pptx_document(tmp_path, pptx_text):
    document = pptx.Presentation()
    document.slides.add_slide(document.slide_layouts[0])
    document.slides[0].shapes.title.text = pptx_text
    document.slides.add_slide(document.slide_layouts[0])
    document.slides[1].shapes.add_textbox(0, 0, 100, 100).text = pptx_text
    path = tmp_path / "test_document.pptx"
    document.save(path)
    return LocalDocument.from_path(path)


def test_pptx(tmp_path):
    pptx_text = "ragna is neat!"
    tmp_pptx_document = get_pptx_document(tmp_path, pptx_text)
    assert isinstance(tmp_pptx_document.handler, PptxDocumentHandler)
    pages = list(tmp_pptx_document.extract_pages())
    assert len(pages) == 2
    for page in pages:
        assert page.text == pptx_text
