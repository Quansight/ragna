from __future__ import annotations

import abc
import io
import uuid
from functools import cached_property
from pathlib import Path
from typing import Any, AsyncIterator, Iterator, Optional, Type, TypeVar, Union

import aiofiles
from pydantic import BaseModel

import ragna

from ._utils import PackageRequirement, RagnaException, Requirement, RequirementsMixin


class Document(RequirementsMixin, abc.ABC):
    """Abstract base class for all documents."""

    def __init__(
        self,
        *,
        id: Optional[uuid.UUID] = None,
        name: str,
        metadata: dict[str, Any],
        handler: Optional[DocumentHandler] = None,
    ):
        self.id = id or uuid.uuid4()
        self.name = name
        self.metadata = metadata
        self.handler = handler or self.get_handler(name)

    @staticmethod
    def supported_suffixes() -> set[str]:
        """
        Returns:
            Suffixes, i.e. `".txt"`, that can be handled by the builtin
                [ragna.core.DocumentHandler][]s.
        """
        return set(DOCUMENT_HANDLERS.keys())

    @staticmethod
    def get_handler(name: str) -> DocumentHandler:
        """Get a document handler based on a suffix.

        Args:
            name: Name of the document.
        """
        handler = DOCUMENT_HANDLERS.get(Path(name).suffix)
        if handler is None:
            raise RagnaException

        return handler

    @abc.abstractmethod
    def read(self) -> bytes: ...

    def extract_pages(self) -> Iterator[Page]:
        yield from self.handler.extract_pages(self)


class LocalDocument(Document):
    """Document class for files on the local file system.

    !!! tip

        This object is usually not instantiated manually, but rather through
        [ragna.core.LocalDocument.from_path][].
    """

    def __init__(
        self,
        *,
        id: Optional[uuid.UUID] = None,
        name: str,
        metadata: dict[str, Any],
        handler: Optional[DocumentHandler] = None,
    ):
        super().__init__(id=id, name=name, metadata=metadata, handler=handler)
        if "path" not in self.metadata:
            metadata["path"] = str(ragna.local_root() / "documents" / str(self.id))

    @classmethod
    def from_path(
        cls,
        path: Union[str, Path],
        *,
        id: Optional[uuid.UUID] = None,
        name: Optional[str] = None,
        metadata: Optional[dict[str, Any]] = None,
        handler: Optional[DocumentHandler] = None,
    ) -> LocalDocument:
        """Create a [ragna.core.LocalDocument][] from a path.

        Args:
            path: Local path to the file.
            id: ID of the document. If omitted, one is generated.
            name: Name of the document. If omitted, defaults to the name of the `path`.
            metadata: Optional metadata of the document.
            handler: Document handler. If omitted, a builtin handler is selected based
                on the suffix of the `path`.

        Raises:
            RagnaException: If `metadata` is passed and contains a `"path"` key.
        """
        if metadata is None:
            metadata = {}
        elif "path" in metadata:
            raise RagnaException(
                "The metadata already includes a 'path' key. "
                "Did you mean to instantiate the class directly?"
            )

        path = Path(path).expanduser().resolve()
        if name is None:
            name = path.name
        metadata["path"] = str(path)
        metadata["extension"] = "".join(path.suffixes)
        metadata["size"] = path.stat().st_size

        return cls(id=id, name=name, metadata=metadata, handler=handler)

    @cached_property
    def path(self) -> Path:
        return Path(self.metadata["path"])

    async def _write(self, stream: AsyncIterator[bytes]) -> None:
        if self.path.exists():
            raise RagnaException(
                "File already exists", path=self.path, http_detail=RagnaException.EVENT
            )

        async with aiofiles.open(self.path, "wb") as file:
            async for content in stream:
                await file.write(content)

    def read(self) -> bytes:
        if not self.path.is_file():
            raise RagnaException(
                "File does not exist", path=self.path, http_detail=RagnaException.EVENT
            )

        with open(self.path, "rb") as file:
            return file.read()


class Page(BaseModel):
    """Dataclass for pages of a document

    Attributes:
        text: Text included in the page.
        number: Page number.
    """

    text: str
    number: Optional[int] = None


class DocumentHandler(RequirementsMixin, abc.ABC):
    """Base class for all document handlers."""

    @classmethod
    @abc.abstractmethod
    def supported_suffixes(cls) -> list[str]:
        """
        Returns:
            Suffixes supported by this document handler.
        """
        pass

    @abc.abstractmethod
    def extract_pages(self, document: Document) -> Iterator[Page]:
        """Extract pages from a document.

        Args:
            document: Document to extract pages from.

        Returns:
            Extracted pages.
        """
        ...


T = TypeVar("T", bound=DocumentHandler)


class DocumentHandlerRegistry(dict[str, DocumentHandler]):
    def load_if_available(self, cls: Type[T]) -> Type[T]:
        if cls.is_available():
            instance = cls()
            for suffix in cls.supported_suffixes():
                self[suffix] = instance

        return cls


DOCUMENT_HANDLERS = DocumentHandlerRegistry()


@DOCUMENT_HANDLERS.load_if_available
class PlainTextDocumentHandler(DocumentHandler):
    """Document handler for plain-text documents.
    Currently supports `.txt` and `.md` extensions.
    """

    @classmethod
    def supported_suffixes(cls) -> list[str]:
        return [".txt", ".md"]

    def extract_pages(self, document: Document) -> Iterator[Page]:
        yield Page(text=document.read().decode())


@DOCUMENT_HANDLERS.load_if_available
class PdfDocumentHandler(DocumentHandler):
    """Document handler for `.pdf` documents.

    !!! info "Package requirements"

        - [`pymupdf`](https://pymupdf.readthedocs.io/en/latest/)
    """

    @classmethod
    def requirements(cls) -> list[Requirement]:
        return [PackageRequirement("pymupdf<=1.24.10,>=1.23.6")]

    @classmethod
    def supported_suffixes(cls) -> list[str]:
        # TODO: pymudpdf supports a lot more formats, while .pdf is by far the most
        #  prominent. Should we expose the others here as well?
        return [".pdf"]

    def extract_pages(self, document: Document) -> Iterator[Page]:
        import fitz

        with fitz.Document(
            stream=document.read(), filetype=Path(document.name).suffix
        ) as document:
            for number, page in enumerate(document, 1):
                yield Page(text=page.get_text(sort=True), number=number)


@DOCUMENT_HANDLERS.load_if_available
class DocxDocumentHandler(DocumentHandler):
    """Document handler for `.docx` documents.

    !!! note

        This does *not* extract text from headers or footers.

    !!! info "Package requirements"

        - [`python-docx`](https://github.com/python-openxml/python-docx)
    """

    @classmethod
    def requirements(cls) -> list[Requirement]:
        return [PackageRequirement("python-docx")]

    @classmethod
    def supported_suffixes(cls) -> list[str]:
        return [".docx"]

    def extract_pages(self, document: Document) -> Iterator[Page]:
        import docx

        document_docx = docx.Document(io.BytesIO(document.read()))
        for paragraph in document_docx.paragraphs:
            text = paragraph.text
            if len(text) > 0:
                yield Page(text=text)


@DOCUMENT_HANDLERS.load_if_available
class PptxDocumentHandler(DocumentHandler):
    """Document handler for `.pptx` documents.

    !!! info "Package requirements"

        - [`python-pptx`](https://github.com/scanny/python-pptx)
    """

    @classmethod
    def requirements(cls) -> list[Requirement]:
        return [PackageRequirement("python-pptx")]

    @classmethod
    def supported_suffixes(cls) -> list[str]:
        return [".pptx"]

    def extract_pages(self, document: Document) -> Iterator[Page]:
        import pptx

        document_pptx = pptx.Presentation(io.BytesIO(document.read()))
        for number, slide in enumerate(document_pptx.slides, 1):
            text = "\n\n".join(
                shape.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text
            )
            yield Page(text=text, number=number)
